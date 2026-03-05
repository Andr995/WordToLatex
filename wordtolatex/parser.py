"""
Modulo per il parsing dei documenti Word, ODT, RTF.
Estrae testo, formattazione, immagini, tabelle, liste, intestazioni.
"""

import os
import json
import re
import zipfile
import subprocess
import tempfile
import shutil
from dataclasses import dataclass, field
from enum import Enum, auto
from typing import Optional
from pathlib import Path

try:
    from docx import Document as DocxDocument
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.enum.style import WD_STYLE_TYPE
    from docx.oxml.ns import qn
    HAS_DOCX = True
except ImportError:
    HAS_DOCX = False

try:
    from odf.opendocument import load as odf_load
    from odf import text as odf_text
    from odf import draw as odf_draw
    from odf import table as odf_table
    from odf.text import P, H, Span, List, ListItem, S
    from odf.element import Element as OdfElement
    HAS_ODF = True
except ImportError:
    HAS_ODF = False


class ElementType(Enum):
    HEADING = auto()
    PARAGRAPH = auto()
    IMAGE = auto()
    TABLE = auto()
    LIST_ITEM = auto()
    PAGE_BREAK = auto()
    HORIZONTAL_RULE = auto()
    CODE_BLOCK = auto()
    FOOTNOTE = auto()


class TextStyle(Enum):
    BOLD = auto()
    ITALIC = auto()
    UNDERLINE = auto()
    STRIKETHROUGH = auto()
    SUPERSCRIPT = auto()
    SUBSCRIPT = auto()
    MONOSPACE = auto()
    SMALL_CAPS = auto()


class ListType(Enum):
    BULLET = auto()
    NUMBERED = auto()


class Alignment(Enum):
    LEFT = auto()
    CENTER = auto()
    RIGHT = auto()
    JUSTIFY = auto()


@dataclass
class TextRun:
    """Un frammento di testo con la sua formattazione."""
    text: str
    styles: set = field(default_factory=set)
    font_size: Optional[float] = None
    font_color: Optional[str] = None
    highlight_color: Optional[str] = None
    hyperlink: Optional[str] = None


@dataclass
class TableCell:
    """Una cella di tabella."""
    runs: list = field(default_factory=list)
    rowspan: int = 1
    colspan: int = 1


@dataclass
class DocumentElement:
    """Elemento generico del documento."""
    element_type: ElementType
    runs: list = field(default_factory=list)        # List[TextRun]
    level: int = 1                                    # per heading/lista
    alignment: Alignment = Alignment.LEFT
    # Per immagini
    image_path: Optional[str] = None
    image_width: Optional[float] = None               # cm
    image_height: Optional[float] = None               # cm
    image_caption: Optional[str] = None
    # Per tabelle
    table_rows: list = field(default_factory=list)    # List[List[TableCell]]
    table_header_rows: int = 0
    # Per liste
    list_type: ListType = ListType.BULLET
    list_depth: int = 0
    # Per footnote
    footnote_text: Optional[str] = None
    # Indentazione
    indent_level: int = 0


class DocumentParser:
    """Parser principale per documenti Word, ODT, RTF."""

    SUPPORTED_EXTENSIONS = {
        '.docx', '.doc', '.odt', '.rtf',
        '.pptx', '.html', '.htm', '.epub',
        '.txt', '.md', '.ipynb',
    }

    def __init__(self, file_path: str, image_output_dir: str = None):
        self.file_path = Path(file_path)
        self.image_output_dir = Path(image_output_dir) if image_output_dir else None
        self.elements: list = []
        self.metadata: dict = {}
        self._image_counter = 0

        if not self.file_path.exists():
            raise FileNotFoundError(f"File non trovato: {self.file_path}")

        ext = self.file_path.suffix.lower()
        if ext not in self.SUPPORTED_EXTENSIONS:
            raise ValueError(
                f"Formato non supportato: {ext}. "
                f"Formati supportati: {', '.join(self.SUPPORTED_EXTENSIONS)}"
            )

    def parse(self) -> list:
        """Parsa il documento e restituisce gli elementi."""
        ext = self.file_path.suffix.lower()

        if ext == '.docx':
            self._parse_docx()
        elif ext == '.doc':
            self._parse_doc()
        elif ext == '.odt':
            self._parse_odt()
        elif ext == '.rtf':
            self._parse_rtf()
        elif ext == '.pptx':
            self._parse_pptx()
        elif ext in ('.html', '.htm'):
            self._parse_html()
        elif ext == '.epub':
            self._parse_epub()
        elif ext == '.txt':
            self._parse_txt()
        elif ext == '.md':
            self._parse_markdown()
        elif ext == '.ipynb':
            self._parse_ipynb()

        return self.elements

    # =========================================================================
    # DOCX PARSING
    # =========================================================================

    def _parse_docx(self):
        """Parsa un file .docx usando python-docx."""
        if not HAS_DOCX:
            raise ImportError(
                "python-docx non installato. Installa con: pip install python-docx"
            )

        doc = DocxDocument(str(self.file_path))
        self._extract_docx_metadata(doc)
        self._extract_docx_images(doc)

        # Itera su tutti gli elementi del body mantenendo l'ordine
        body = doc.element.body
        for child in body:
            tag = child.tag.split('}')[-1] if '}' in child.tag else child.tag

            if tag == 'p':
                self._parse_docx_paragraph(child, doc)
            elif tag == 'tbl':
                self._parse_docx_table(child, doc)
            elif tag == 'sectPr':
                pass  # Sezione, ignorata

    def _extract_docx_metadata(self, doc):
        """Estrae metadati dal documento docx."""
        props = doc.core_properties
        self.metadata = {
            'title': props.title or '',
            'author': props.author or '',
            'subject': props.subject or '',
            'created': str(props.created) if props.created else '',
            'modified': str(props.modified) if props.modified else '',
        }

    def _extract_docx_images(self, doc):
        """Estrae tutte le immagini dal file docx."""
        if not self.image_output_dir:
            return

        self.image_output_dir.mkdir(parents=True, exist_ok=True)

        for rel in doc.part.rels.values():
            if "image" in rel.reltype:
                try:
                    image_data = rel.target_part.blob
                    content_type = rel.target_part.content_type
                    ext_map = {
                        'image/png': '.png',
                        'image/jpeg': '.jpg',
                        'image/gif': '.gif',
                        'image/bmp': '.bmp',
                        'image/tiff': '.tiff',
                        'image/svg+xml': '.svg',
                        'image/x-emf': '.emf',
                        'image/x-wmf': '.wmf',
                    }
                    ext = ext_map.get(content_type, '.png')
                    # Usa il nome originale del file immagine
                    img_name = rel.target_ref.split('/')[-1]
                    if not img_name:
                        self._image_counter += 1
                        img_name = f"image_{self._image_counter}{ext}"

                    img_path = self.image_output_dir / img_name
                    with open(img_path, 'wb') as f:
                        f.write(image_data)
                except Exception:
                    pass

    def _parse_docx_paragraph(self, para_elem, doc):
        """Parsa un paragrafo docx."""
        from docx.text.paragraph import Paragraph
        para = Paragraph(para_elem, doc)

        # Controlla se è un page break
        if self._docx_has_page_break(para):
            self.elements.append(DocumentElement(element_type=ElementType.PAGE_BREAK))

        # Controlla stile
        style_name = para.style.name if para.style else ''
        style_lower = style_name.lower()

        # Heading
        if style_lower.startswith('heading') or style_lower.startswith('titre'):
            level = 1
            # Prova a estrarre il livello dal nome stile
            for char in style_name:
                if char.isdigit():
                    level = int(char)
                    break
            runs = self._docx_extract_runs(para)
            if runs:
                elem = DocumentElement(
                    element_type=ElementType.HEADING,
                    runs=runs,
                    level=level,
                    alignment=self._docx_get_alignment(para),
                )
                self.elements.append(elem)
            return

        # Titolo (Title)
        if 'title' in style_lower or 'titolo' in style_lower:
            runs = self._docx_extract_runs(para)
            if runs:
                elem = DocumentElement(
                    element_type=ElementType.HEADING,
                    runs=runs,
                    level=0,  # Titolo del documento
                    alignment=Alignment.CENTER,
                )
                self.elements.append(elem)
            return

        # Subtitle
        if 'subtitle' in style_lower or 'sottotitolo' in style_lower:
            runs = self._docx_extract_runs(para)
            if runs:
                elem = DocumentElement(
                    element_type=ElementType.HEADING,
                    runs=runs,
                    level=0,  # Anche il sottotitolo è trattato come titolo
                    alignment=Alignment.CENTER,
                )
                self.elements.append(elem)
            return

        # Lista
        if self._docx_is_list(para):
            runs = self._docx_extract_runs(para)
            if runs:
                list_type, depth = self._docx_get_list_info(para)
                elem = DocumentElement(
                    element_type=ElementType.LIST_ITEM,
                    runs=runs,
                    list_type=list_type,
                    list_depth=depth,
                    alignment=self._docx_get_alignment(para),
                )
                self.elements.append(elem)
            return

        # Controlla immagini inline
        images = self._docx_extract_inline_images(para, doc)
        for img_elem in images:
            self.elements.append(img_elem)

        # Paragrafo normale
        runs = self._docx_extract_runs(para)
        if runs or not images:
            indent = self._docx_get_indent_level(para)
            elem = DocumentElement(
                element_type=ElementType.PARAGRAPH,
                runs=runs,
                alignment=self._docx_get_alignment(para),
                indent_level=indent,
            )
            self.elements.append(elem)

    def _docx_extract_runs(self, para) -> list:
        """Estrae i TextRun da un paragrafo docx."""
        text_runs = []
        for run in para.runs:
            styles = set()
            if run.bold:
                styles.add(TextStyle.BOLD)
            if run.italic:
                styles.add(TextStyle.ITALIC)
            if run.underline:
                styles.add(TextStyle.UNDERLINE)
            if run.font.strike:
                styles.add(TextStyle.STRIKETHROUGH)
            if run.font.superscript:
                styles.add(TextStyle.SUPERSCRIPT)
            if run.font.subscript:
                styles.add(TextStyle.SUBSCRIPT)
            if run.font.small_caps:
                styles.add(TextStyle.SMALL_CAPS)

            # Font monospace
            font_name = run.font.name or ''
            if any(m in font_name.lower() for m in ['courier', 'mono', 'consolas', 'menlo']):
                styles.add(TextStyle.MONOSPACE)

            # Dimensione font in pt
            font_size = None
            if run.font.size:
                font_size = run.font.size.pt

            # Colore font
            font_color = None
            if run.font.color and run.font.color.rgb:
                font_color = str(run.font.color.rgb)

            text_runs.append(TextRun(
                text=run.text,
                styles=styles,
                font_size=font_size,
                font_color=font_color,
            ))

        # Controlla hyperlink
        for link in para._element.findall(qn('w:hyperlink')):
            r_id = link.get(qn('r:id'))
            url = None
            if r_id:
                try:
                    rel = para.part.rels.get(r_id)
                    if rel:
                        url = rel.target_ref
                except Exception:
                    pass

            link_text = ''.join(
                node.text or '' for node in link.findall(qn('w:r') + '/' + qn('w:t'))
            )
            if link_text:
                text_runs.append(TextRun(
                    text=link_text,
                    hyperlink=url,
                    styles={TextStyle.UNDERLINE},
                ))

        return text_runs

    def _docx_has_page_break(self, para) -> bool:
        """Verifica se il paragrafo contiene un page break."""
        for run in para.runs:
            br_elements = run._element.findall(qn('w:br'))
            for br in br_elements:
                if br.get(qn('w:type')) == 'page':
                    return True
        return False

    def _docx_get_alignment(self, para) -> Alignment:
        """Ottiene l'allineamento del paragrafo."""
        try:
            alignment = para.alignment
            if alignment == WD_ALIGN_PARAGRAPH.CENTER:
                return Alignment.CENTER
            elif alignment == WD_ALIGN_PARAGRAPH.RIGHT:
                return Alignment.RIGHT
            elif alignment == WD_ALIGN_PARAGRAPH.JUSTIFY:
                return Alignment.JUSTIFY
        except Exception:
            pass
        return Alignment.LEFT

    def _docx_is_list(self, para) -> bool:
        """Verifica se il paragrafo è un elemento di lista."""
        pPr = para._element.find(qn('w:pPr'))
        if pPr is not None:
            numPr = pPr.find(qn('w:numPr'))
            if numPr is not None:
                return True

        style_name = (para.style.name or '').lower()
        if 'list' in style_name or 'elenco' in style_name:
            return True

        return False

    def _docx_get_list_info(self, para):
        """Ottiene tipo e profondità della lista."""
        depth = 0
        list_type = ListType.BULLET

        pPr = para._element.find(qn('w:pPr'))
        if pPr is not None:
            numPr = pPr.find(qn('w:numPr'))
            if numPr is not None:
                ilvl = numPr.find(qn('w:ilvl'))
                if ilvl is not None:
                    depth = int(ilvl.get(qn('w:val'), '0'))

                numId = numPr.find(qn('w:numId'))
                if numId is not None:
                    num_val = numId.get(qn('w:val'), '0')
                    # Euristica: numId dispari tende ad essere numerato
                    if int(num_val) % 2 == 0:
                        list_type = ListType.NUMBERED

        style_name = (para.style.name or '').lower()
        if 'number' in style_name or 'numer' in style_name:
            list_type = ListType.NUMBERED

        return list_type, depth

    def _docx_get_indent_level(self, para) -> int:
        """Calcola il livello di indentazione."""
        try:
            pPr = para._element.find(qn('w:pPr'))
            if pPr is not None:
                ind = pPr.find(qn('w:ind'))
                if ind is not None:
                    left = ind.get(qn('w:left'), '0')
                    # Converti twips in livello (720 twips = 1 livello)
                    return int(int(left) / 720)
        except Exception:
            pass
        return 0

    def _docx_extract_inline_images(self, para, doc) -> list:
        """Estrae immagini inline da un paragrafo."""
        images = []
        if not self.image_output_dir:
            return images

        # Cerca drawing elements
        for drawing in para._element.findall('.//' + qn('w:drawing')):
            # Cerca inline images
            for inline in drawing.findall('.//' + qn('wp:inline')):
                img_elem = self._docx_process_image_element(inline, doc)
                if img_elem:
                    images.append(img_elem)

            # Cerca anchor images
            for anchor in drawing.findall('.//' + qn('wp:anchor')):
                img_elem = self._docx_process_image_element(anchor, doc)
                if img_elem:
                    images.append(img_elem)

        return images

    def _docx_process_image_element(self, element, doc) -> Optional[DocumentElement]:
        """Processa un elemento immagine docx."""
        try:
            # Prova a trovare il blip (riferimento all'immagine)
            blip = element.find('.//' + qn('a:blip'))
            if blip is None:
                return None

            r_id = blip.get(qn('r:embed'))
            if not r_id:
                return None

            rel = doc.part.rels.get(r_id)
            if rel is None:
                return None

            img_name = rel.target_ref.split('/')[-1]
            img_path = self.image_output_dir / img_name

            # Estrai dimensioni (EMU -> cm, 1 cm = 914400 EMU)
            width_cm = None
            height_cm = None
            extent = element.find(qn('wp:extent'))
            if extent is not None:
                cx = extent.get('cx')
                cy = extent.get('cy')
                if cx:
                    width_cm = int(cx) / 914400
                if cy:
                    height_cm = int(cy) / 914400

            if img_path.exists():
                return DocumentElement(
                    element_type=ElementType.IMAGE,
                    image_path=str(img_path),
                    image_width=width_cm,
                    image_height=height_cm,
                )
        except Exception:
            pass

        return None

    def _parse_docx_table(self, tbl_elem, doc):
        """Parsa una tabella docx."""
        from docx.table import Table
        table = Table(tbl_elem, doc)

        rows = []
        for row in table.rows:
            cells = []
            for cell in row.cells:
                cell_runs = []
                for para_idx, para in enumerate(cell.paragraphs):
                    for run in para.runs:
                        styles = set()
                        if run.bold:
                            styles.add(TextStyle.BOLD)
                        if run.italic:
                            styles.add(TextStyle.ITALIC)
                        # Strip trailing newlines dal testo delle run
                        text = run.text.rstrip('\n').rstrip('\r')
                        if text:
                            cell_runs.append(TextRun(text=text, styles=styles))
                    # Aggiungi separatore tra paragrafi nella cella
                    if para_idx < len(cell.paragraphs) - 1 and cell_runs:
                        cell_runs.append(TextRun(text=' '))

                tc = TableCell(runs=cell_runs)

                # Gestisci merge
                tc_elem = cell._tc
                grid_span = tc_elem.find(qn('w:tcPr'))
                if grid_span is not None:
                    gs = grid_span.find(qn('w:gridSpan'))
                    if gs is not None:
                        tc.colspan = int(gs.get(qn('w:val'), '1'))
                    vm = grid_span.find(qn('w:vMerge'))
                    if vm is not None:
                        val = vm.get(qn('w:val'), '')
                        if val != 'restart':
                            tc.rowspan = 0  # Cella merged (continua)

                cells.append(tc)
            rows.append(cells)

        elem = DocumentElement(
            element_type=ElementType.TABLE,
            table_rows=rows,
            table_header_rows=1,  # Assume prima riga come header
        )
        self.elements.append(elem)

    # =========================================================================
    # DOC PARSING (via conversione a docx tramite LibreOffice)
    # =========================================================================

    def _parse_doc(self):
        """Parsa un file .doc convertendolo prima in .docx via LibreOffice."""
        docx_path = self._convert_with_libreoffice('.docx')
        if docx_path:
            original_path = self.file_path
            self.file_path = Path(docx_path)
            self._parse_docx()
            self.file_path = original_path
            # Pulizia
            try:
                os.unlink(docx_path)
            except OSError:
                pass

    # =========================================================================
    # ODT PARSING
    # =========================================================================

    def _parse_odt(self):
        """Parsa un file .odt usando odfpy."""
        if not HAS_ODF:
            # Fallback: converti con LibreOffice
            docx_path = self._convert_with_libreoffice('.docx')
            if docx_path:
                original_path = self.file_path
                self.file_path = Path(docx_path)
                self._parse_docx()
                self.file_path = original_path
                try:
                    os.unlink(docx_path)
                except OSError:
                    pass
            return

        doc = odf_load(str(self.file_path))
        self._extract_odt_images(doc)
        body = doc.body

        for elem in body.childNodes:
            tag = elem.qname[1] if isinstance(elem.qname, tuple) else str(elem.qname)

            if tag == 'p':
                self._parse_odt_paragraph(elem)
            elif tag == 'h':
                self._parse_odt_heading(elem)
            elif tag == 'table':
                self._parse_odt_table(elem)
            elif tag == 'list':
                self._parse_odt_list(elem, depth=0)

    def _extract_odt_images(self, doc):
        """Estrae immagini da un file ODT."""
        if not self.image_output_dir:
            return
        self.image_output_dir.mkdir(parents=True, exist_ok=True)

        # ODT è uno zip, estrai le immagini dalla cartella Pictures/
        try:
            with zipfile.ZipFile(str(self.file_path), 'r') as zf:
                for name in zf.namelist():
                    if name.startswith('Pictures/'):
                        img_name = os.path.basename(name)
                        if img_name:
                            img_path = self.image_output_dir / img_name
                            with zf.open(name) as src, open(img_path, 'wb') as dst:
                                dst.write(src.read())
        except Exception:
            pass

    def _parse_odt_paragraph(self, elem):
        """Parsa un paragrafo ODT."""
        runs = self._odt_extract_text_runs(elem)
        self.elements.append(DocumentElement(
            element_type=ElementType.PARAGRAPH,
            runs=runs,
        ))

    def _parse_odt_heading(self, elem):
        """Parsa un heading ODT."""
        level = 1
        try:
            outline_level = elem.getAttribute('outlinelevel')
            if outline_level:
                level = int(outline_level)
        except (ValueError, AttributeError):
            pass

        runs = self._odt_extract_text_runs(elem)
        self.elements.append(DocumentElement(
            element_type=ElementType.HEADING,
            runs=runs,
            level=level,
        ))

    def _parse_odt_table(self, elem):
        """Parsa una tabella ODT."""
        rows = []
        for child in elem.childNodes:
            tag = child.qname[1] if isinstance(child.qname, tuple) else str(child.qname)
            if tag == 'table-row':
                row_cells = []
                for cell_elem in child.childNodes:
                    cell_tag = cell_elem.qname[1] if isinstance(cell_elem.qname, tuple) else str(cell_elem.qname)
                    if cell_tag == 'table-cell':
                        runs = self._odt_extract_text_runs(cell_elem)
                        row_cells.append(TableCell(runs=runs))
                if row_cells:
                    rows.append(row_cells)

        if rows:
            self.elements.append(DocumentElement(
                element_type=ElementType.TABLE,
                table_rows=rows,
                table_header_rows=1,
            ))

    def _parse_odt_list(self, elem, depth=0):
        """Parsa una lista ODT."""
        for child in elem.childNodes:
            tag = child.qname[1] if isinstance(child.qname, tuple) else str(child.qname)
            if tag == 'list-item':
                for subchild in child.childNodes:
                    subtag = subchild.qname[1] if isinstance(subchild.qname, tuple) else str(subchild.qname)
                    if subtag == 'p':
                        runs = self._odt_extract_text_runs(subchild)
                        self.elements.append(DocumentElement(
                            element_type=ElementType.LIST_ITEM,
                            runs=runs,
                            list_type=ListType.BULLET,
                            list_depth=depth,
                        ))
                    elif subtag == 'list':
                        self._parse_odt_list(subchild, depth + 1)

    def _odt_extract_text_runs(self, elem) -> list:
        """Estrae TextRun da un elemento ODT."""
        runs = []
        self._odt_walk_text(elem, runs, set())
        return runs

    def _odt_walk_text(self, node, runs, current_styles):
        """Attraversa ricorsivamente i nodi ODT per estrarre testo."""
        if hasattr(node, 'data') and node.data:
            # Nodo testo
            runs.append(TextRun(text=node.data, styles=set(current_styles)))
            return

        tag = ''
        if hasattr(node, 'qname'):
            tag = node.qname[1] if isinstance(node.qname, tuple) else str(node.qname)

        new_styles = set(current_styles)

        # Gestisci lo stile span
        if tag == 'span':
            style_name = ''
            try:
                style_name = (node.getAttribute('stylename') or '').lower()
            except Exception:
                pass
            if 'bold' in style_name or 'grassetto' in style_name:
                new_styles.add(TextStyle.BOLD)
            if 'italic' in style_name or 'corsivo' in style_name:
                new_styles.add(TextStyle.ITALIC)

        # Gestisci spazi multipli
        if tag == 's':
            count = 1
            try:
                c = node.getAttribute('c')
                if c:
                    count = int(c)
            except Exception:
                pass
            runs.append(TextRun(text=' ' * count, styles=set(current_styles)))
            return

        # Tab
        if tag == 'tab':
            runs.append(TextRun(text='\t', styles=set(current_styles)))
            return

        # Line break
        if tag == 'line-break':
            runs.append(TextRun(text='\n', styles=set(current_styles)))
            return

        # Immagine
        if tag == 'frame':
            for child in node.childNodes:
                child_tag = child.qname[1] if isinstance(child.qname, tuple) else str(child.qname)
                if child_tag == 'image':
                    href = child.getAttribute('href') or ''
                    if href:
                        img_name = os.path.basename(href)
                        if self.image_output_dir:
                            img_path = self.image_output_dir / img_name
                            if img_path.exists():
                                self.elements.append(DocumentElement(
                                    element_type=ElementType.IMAGE,
                                    image_path=str(img_path),
                                ))
            return

        # Ricorsione sui figli
        if hasattr(node, 'childNodes'):
            for child in node.childNodes:
                self._odt_walk_text(child, runs, new_styles)

    # =========================================================================
    # RTF PARSING (via conversione)
    # =========================================================================

    def _parse_rtf(self):
        """Parsa un file .rtf convertendolo a .docx via LibreOffice."""
        docx_path = self._convert_with_libreoffice('.docx')
        if docx_path:
            original_path = self.file_path
            self.file_path = Path(docx_path)
            self._parse_docx()
            self.file_path = original_path
            try:
                os.unlink(docx_path)
            except OSError:
                pass

    # =========================================================================
    # UTILITY
    # =========================================================================

    def _convert_with_libreoffice(self, target_ext: str) -> Optional[str]:
        """Converte un documento usando LibreOffice in modalità headless."""
        # Cerca LibreOffice
        lo_cmd = None
        for cmd in ['libreoffice', 'soffice', '/usr/bin/libreoffice',
                     '/usr/bin/soffice', '/snap/bin/libreoffice']:
            if shutil.which(cmd):
                lo_cmd = cmd
                break

        if not lo_cmd:
            raise RuntimeError(
                "LibreOffice non trovato. Installalo per convertire file .doc/.rtf.\n"
                "Su Ubuntu/Debian: sudo apt install libreoffice\n"
                "Su Fedora: sudo dnf install libreoffice\n"
                "Su macOS: brew install --cask libreoffice"
            )

        with tempfile.TemporaryDirectory() as tmpdir:
            # Copia il file sorgente nel tmpdir
            src_copy = os.path.join(tmpdir, self.file_path.name)
            shutil.copy2(str(self.file_path), src_copy)

            fmt = 'docx' if target_ext == '.docx' else target_ext.lstrip('.')

            result = subprocess.run(
                [lo_cmd, '--headless', '--convert-to', fmt,
                 '--outdir', tmpdir, src_copy],
                capture_output=True, text=True, timeout=120,
            )

            if result.returncode != 0:
                raise RuntimeError(
                    f"Errore nella conversione con LibreOffice: {result.stderr}"
                )

            # Trova il file convertito
            base_name = self.file_path.stem + target_ext
            converted = os.path.join(tmpdir, base_name)

            if os.path.exists(converted):
                # Copia in una posizione persistente
                final_path = os.path.join(
                    str(self.image_output_dir or tempfile.gettempdir()),
                    base_name,
                )
                shutil.copy2(converted, final_path)
                return final_path

        return None

    # =========================================================================
    # PPTX PARSING
    # =========================================================================

    def _parse_pptx(self):
        """Parsa un file .pptx usando python-pptx."""
        try:
            from pptx import Presentation
            from pptx.util import Inches, Emu
        except ImportError:
            raise ImportError(
                "python-pptx non installato. Installa con: pip install python-pptx"
            )

        prs = Presentation(str(self.file_path))
        self.metadata['title'] = prs.core_properties.title or self.file_path.stem

        for slide_num, slide in enumerate(prs.slides, 1):
            # Intestazione per ogni slide
            if slide_num > 1:
                self.elements.append(
                    DocumentElement(element_type=ElementType.PAGE_BREAK)
                )

            # Cerca il titolo della slide
            title_text = None
            if slide.shapes.title and slide.shapes.title.has_text_frame:
                title_text = slide.shapes.title.text_frame.text.strip()

            if title_text:
                self.elements.append(DocumentElement(
                    element_type=ElementType.HEADING,
                    runs=[TextRun(text=title_text)],
                    level=2,
                ))

            # Processa tutte le shape nella slide
            for shape in slide.shapes:
                # Salta il titolo (già gestito)
                if shape == slide.shapes.title:
                    continue

                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        runs = []
                        for run in para.runs:
                            styles = set()
                            if run.font.bold:
                                styles.add(TextStyle.BOLD)
                            if run.font.italic:
                                styles.add(TextStyle.ITALIC)
                            if run.font.underline:
                                styles.add(TextStyle.UNDERLINE)
                            if run.text.strip():
                                runs.append(TextRun(
                                    text=run.text,
                                    styles=styles,
                                ))
                        if runs:
                            self.elements.append(DocumentElement(
                                element_type=ElementType.PARAGRAPH,
                                runs=runs,
                            ))

                if shape.has_table:
                    tbl = shape.table
                    rows = []
                    for row in tbl.rows:
                        row_cells = []
                        for cell in row.cells:
                            cell_runs = [TextRun(text=cell.text.strip())]
                            row_cells.append(TableCell(runs=cell_runs))
                        rows.append(row_cells)
                    if rows:
                        self.elements.append(DocumentElement(
                            element_type=ElementType.TABLE,
                            table_rows=rows,
                            table_header_rows=1,
                        ))

                # Immagini dalla presentazione
                if shape.shape_type == 13:  # Picture
                    if self.image_output_dir:
                        try:
                            image = shape.image
                            img_ext = image.content_type.split('/')[-1]
                            if img_ext == 'jpeg':
                                img_ext = 'jpg'
                            self._image_counter += 1
                            img_name = f"slide{slide_num}_img{self._image_counter}.{img_ext}"
                            img_path = self.image_output_dir / img_name
                            with open(img_path, 'wb') as f:
                                f.write(image.blob)

                            width_cm = None
                            if shape.width:
                                width_cm = shape.width / 914400

                            self.elements.append(DocumentElement(
                                element_type=ElementType.IMAGE,
                                image_path=str(img_path),
                                image_width=width_cm,
                            ))
                        except Exception:
                            pass

    # =========================================================================
    # HTML PARSING
    # =========================================================================

    def _parse_html(self):
        """Parsa un file HTML."""
        from html.parser import HTMLParser

        with open(str(self.file_path), 'r', encoding='utf-8', errors='replace') as f:
            html_content = f.read()

        # Rimuovi script e style
        html_content = re.sub(r'<script[^>]*>.*?</script>', '', html_content,
                              flags=re.DOTALL | re.IGNORECASE)
        html_content = re.sub(r'<style[^>]*>.*?</style>', '', html_content,
                              flags=re.DOTALL | re.IGNORECASE)

        class _HTMLContentParser(HTMLParser):
            def __init__(self):
                super().__init__()
                self.elements = []
                self.current_runs = []
                self.current_styles = set()
                self.in_heading = 0
                self.in_list = False
                self.list_type = ListType.BULLET
                self.in_table = False
                self.table_rows = []
                self.current_row = []
                self.current_cell_text = ''
                self.in_pre = False
                self.skip = False

            def handle_starttag(self, tag, attrs):
                tag = tag.lower()
                if tag in ('script', 'style', 'head'):
                    self.skip = True
                    return
                if tag in ('h1', 'h2', 'h3', 'h4', 'h5', 'h6'):
                    self._flush_paragraph()
                    self.in_heading = int(tag[1])
                elif tag == 'b' or tag == 'strong':
                    self.current_styles.add(TextStyle.BOLD)
                elif tag == 'i' or tag == 'em':
                    self.current_styles.add(TextStyle.ITALIC)
                elif tag == 'u':
                    self.current_styles.add(TextStyle.UNDERLINE)
                elif tag in ('s', 'strike', 'del'):
                    self.current_styles.add(TextStyle.STRIKETHROUGH)
                elif tag == 'code' and not self.in_pre:
                    self.current_styles.add(TextStyle.MONOSPACE)
                elif tag == 'pre':
                    self.in_pre = True
                elif tag == 'sup':
                    self.current_styles.add(TextStyle.SUPERSCRIPT)
                elif tag == 'sub':
                    self.current_styles.add(TextStyle.SUBSCRIPT)
                elif tag == 'br':
                    self.current_runs.append(TextRun(text='\n'))
                elif tag == 'hr':
                    self._flush_paragraph()
                    self.elements.append(DocumentElement(
                        element_type=ElementType.HORIZONTAL_RULE
                    ))
                elif tag == 'ul':
                    self._flush_paragraph()
                    self.in_list = True
                    self.list_type = ListType.BULLET
                elif tag == 'ol':
                    self._flush_paragraph()
                    self.in_list = True
                    self.list_type = ListType.NUMBERED
                elif tag == 'li':
                    self.current_runs = []
                elif tag == 'table':
                    self._flush_paragraph()
                    self.in_table = True
                    self.table_rows = []
                elif tag == 'tr':
                    self.current_row = []
                elif tag in ('td', 'th'):
                    self.current_cell_text = ''
                elif tag == 'a':
                    href = dict(attrs).get('href', '')
                    if href:
                        self.current_runs.append(TextRun(
                            text='',  # placeholder, testo sarà nel handle_data
                            hyperlink=href,
                            styles={TextStyle.UNDERLINE},
                        ))
                elif tag == 'img':
                    src = dict(attrs).get('src', '')
                    alt = dict(attrs).get('alt', '')
                    if src:
                        self.elements.append(DocumentElement(
                            element_type=ElementType.IMAGE,
                            image_path=src,
                            image_caption=alt,
                        ))
                elif tag == 'p':
                    self._flush_paragraph()

            def handle_endtag(self, tag):
                tag = tag.lower()
                if tag in ('script', 'style', 'head'):
                    self.skip = False
                    return
                if tag in ('h1', 'h2', 'h3', 'h4', 'h5', 'h6'):
                    level = int(tag[1])
                    if self.current_runs:
                        self.elements.append(DocumentElement(
                            element_type=ElementType.HEADING,
                            runs=list(self.current_runs),
                            level=level,
                        ))
                        self.current_runs = []
                    self.in_heading = 0
                elif tag in ('b', 'strong'):
                    self.current_styles.discard(TextStyle.BOLD)
                elif tag in ('i', 'em'):
                    self.current_styles.discard(TextStyle.ITALIC)
                elif tag == 'u':
                    self.current_styles.discard(TextStyle.UNDERLINE)
                elif tag in ('s', 'strike', 'del'):
                    self.current_styles.discard(TextStyle.STRIKETHROUGH)
                elif tag == 'code' and not self.in_pre:
                    self.current_styles.discard(TextStyle.MONOSPACE)
                elif tag == 'pre':
                    self.in_pre = False
                elif tag == 'sup':
                    self.current_styles.discard(TextStyle.SUPERSCRIPT)
                elif tag == 'sub':
                    self.current_styles.discard(TextStyle.SUBSCRIPT)
                elif tag == 'li':
                    if self.current_runs:
                        self.elements.append(DocumentElement(
                            element_type=ElementType.LIST_ITEM,
                            runs=list(self.current_runs),
                            list_type=self.list_type,
                        ))
                        self.current_runs = []
                elif tag in ('ul', 'ol'):
                    self.in_list = False
                elif tag in ('td', 'th'):
                    self.current_row.append(
                        TableCell(runs=[TextRun(text=self.current_cell_text.strip())])
                    )
                elif tag == 'tr':
                    if self.current_row:
                        self.table_rows.append(self.current_row)
                elif tag == 'table':
                    if self.table_rows:
                        self.elements.append(DocumentElement(
                            element_type=ElementType.TABLE,
                            table_rows=self.table_rows,
                            table_header_rows=1,
                        ))
                    self.in_table = False
                elif tag == 'p':
                    self._flush_paragraph()

            def handle_data(self, data):
                if self.skip:
                    return
                if self.in_table:
                    self.current_cell_text += data
                    return
                text = data
                if not self.in_pre:
                    text = re.sub(r'\s+', ' ', text)
                if text:
                    # Se l'ultimo run era un hyperlink placeholder vuoto, assegna il testo
                    if (self.current_runs and self.current_runs[-1].hyperlink
                            and not self.current_runs[-1].text):
                        self.current_runs[-1] = TextRun(
                            text=text,
                            hyperlink=self.current_runs[-1].hyperlink,
                            styles=self.current_runs[-1].styles,
                        )
                    else:
                        self.current_runs.append(TextRun(
                            text=text,
                            styles=set(self.current_styles),
                        ))

            def _flush_paragraph(self):
                if self.current_runs:
                    full_text = ''.join(r.text for r in self.current_runs).strip()
                    if full_text:
                        self.elements.append(DocumentElement(
                            element_type=ElementType.PARAGRAPH,
                            runs=list(self.current_runs),
                        ))
                    self.current_runs = []

            def finalize(self):
                self._flush_paragraph()

        p = _HTMLContentParser()
        p.feed(html_content)
        p.finalize()
        self.elements = p.elements

        # Prova a estrarre il titolo dall'HTML
        title_match = re.search(r'<title[^>]*>(.*?)</title>',
                                html_content, re.IGNORECASE | re.DOTALL)
        if title_match:
            self.metadata['title'] = title_match.group(1).strip()

    # =========================================================================
    # EPUB PARSING
    # =========================================================================

    def _parse_epub(self):
        """Parsa un file .epub (è uno zip di file HTML)."""
        import zipfile

        if not zipfile.is_zipfile(str(self.file_path)):
            raise ValueError("Il file .epub non è un archivio valido.")

        html_parts = []

        with zipfile.ZipFile(str(self.file_path), 'r') as zf:
            # Estrai immagini
            if self.image_output_dir:
                self.image_output_dir.mkdir(parents=True, exist_ok=True)
                for name in zf.namelist():
                    lower = name.lower()
                    if any(lower.endswith(e) for e in
                           ('.png', '.jpg', '.jpeg', '.gif', '.svg', '.bmp')):
                        img_name = os.path.basename(name)
                        if img_name:
                            img_path = self.image_output_dir / img_name
                            with zf.open(name) as src, open(img_path, 'wb') as dst:
                                dst.write(src.read())

            # Leggi i file HTML/XHTML nell'ordine dell'OPF (spine)
            opf_path = None
            for name in zf.namelist():
                if name.endswith('.opf'):
                    opf_path = name
                    break

            html_files = []
            if opf_path:
                opf_content = zf.read(opf_path).decode('utf-8', errors='replace')
                opf_dir = os.path.dirname(opf_path)

                # Estrai manifest
                manifest = {}
                for match in re.finditer(
                    r'<item\s+[^>]*id="([^"]*)"[^>]*href="([^"]*)"[^>]*/?>',
                    opf_content, re.DOTALL
                ):
                    item_id = match.group(1)
                    href = match.group(2)
                    full_path = os.path.join(opf_dir, href) if opf_dir else href
                    manifest[item_id] = full_path

                # Estrai spine order
                for match in re.finditer(
                    r'<itemref\s+[^>]*idref="([^"]*)"',
                    opf_content
                ):
                    item_id = match.group(1)
                    if item_id in manifest:
                        html_files.append(manifest[item_id])

            # Fallback: tutti i file html/xhtml
            if not html_files:
                html_files = [n for n in zf.namelist()
                              if n.lower().endswith(('.html', '.xhtml', '.htm'))]
                html_files.sort()

            for html_file in html_files:
                try:
                    content = zf.read(html_file).decode('utf-8', errors='replace')
                    html_parts.append(content)
                except Exception:
                    pass

        # Parsa l'HTML combinato
        combined = '\n'.join(html_parts)
        # Salva temporaneamente e usa il parser HTML
        original = self.file_path
        with tempfile.NamedTemporaryFile(mode='w', suffix='.html',
                                         delete=False, encoding='utf-8') as f:
            f.write(combined)
            tmp_html = f.name

        try:
            self.file_path = Path(tmp_html)
            self._parse_html()
            self.file_path = original
        finally:
            try:
                os.unlink(tmp_html)
            except OSError:
                pass

    # =========================================================================
    # TXT PARSING
    # =========================================================================

    def _parse_txt(self):
        """Parsa un file di testo semplice."""
        with open(str(self.file_path), 'r', encoding='utf-8', errors='replace') as f:
            content = f.read()

        self.metadata['title'] = self.file_path.stem

        # Dividi in paragrafi (righe vuote separano i paragrafi)
        paragraphs = re.split(r'\n\s*\n', content)

        for para in paragraphs:
            text = para.strip()
            if not text:
                continue

            # Euristica: linee tutte maiuscole o brevi possono essere titoli
            lines = text.split('\n')
            if len(lines) == 1 and len(text) < 80 and text == text.upper() and len(text) > 3:
                self.elements.append(DocumentElement(
                    element_type=ElementType.HEADING,
                    runs=[TextRun(text=text.title())],
                    level=1,
                ))
            elif len(lines) == 1 and text.startswith('# '):
                # Supporto minimo per heading markdown-like
                self.elements.append(DocumentElement(
                    element_type=ElementType.HEADING,
                    runs=[TextRun(text=text.lstrip('# ').strip())],
                    level=1,
                ))
            else:
                # Preserva le righe interne come un unico paragrafo
                joined = ' '.join(line.strip() for line in lines if line.strip())
                self.elements.append(DocumentElement(
                    element_type=ElementType.PARAGRAPH,
                    runs=[TextRun(text=joined)],
                ))

    # =========================================================================
    # MARKDOWN PARSING
    # =========================================================================

    def _parse_markdown(self):
        """Parsa un file Markdown."""
        with open(str(self.file_path), 'r', encoding='utf-8', errors='replace') as f:
            lines = f.readlines()

        self.metadata['title'] = self.file_path.stem

        i = 0
        in_code_block = False
        code_lines = []
        in_list = False

        while i < len(lines):
            line = lines[i]
            stripped = line.rstrip('\n')

            # Code blocks
            if stripped.startswith('```'):
                if in_code_block:
                    # Fine code block
                    code_text = '\n'.join(code_lines)
                    self.elements.append(DocumentElement(
                        element_type=ElementType.CODE_BLOCK,
                        runs=[TextRun(text=code_text)],
                    ))
                    code_lines = []
                    in_code_block = False
                else:
                    in_code_block = True
                i += 1
                continue

            if in_code_block:
                code_lines.append(stripped)
                i += 1
                continue

            # Righe vuote
            if not stripped.strip():
                in_list = False
                i += 1
                continue

            # Headings
            heading_match = re.match(r'^(#{1,6})\s+(.+)', stripped)
            if heading_match:
                level = len(heading_match.group(1))
                text = heading_match.group(2).strip()
                # Prima heading h1 diventa il titolo
                if level == 1 and not self.metadata.get('title_set'):
                    self.metadata['title'] = text
                    self.metadata['title_set'] = True
                runs = self._md_parse_inline(text)
                self.elements.append(DocumentElement(
                    element_type=ElementType.HEADING,
                    runs=runs,
                    level=level,
                ))
                i += 1
                continue

            # Horizontal rule
            if re.match(r'^[-*_]{3,}\s*$', stripped):
                self.elements.append(DocumentElement(
                    element_type=ElementType.HORIZONTAL_RULE
                ))
                i += 1
                continue

            # Liste puntate
            list_match = re.match(r'^(\s*)[-*+]\s+(.+)', stripped)
            if list_match:
                depth = len(list_match.group(1)) // 2
                text = list_match.group(2).strip()
                runs = self._md_parse_inline(text)
                self.elements.append(DocumentElement(
                    element_type=ElementType.LIST_ITEM,
                    runs=runs,
                    list_type=ListType.BULLET,
                    list_depth=depth,
                ))
                in_list = True
                i += 1
                continue

            # Liste numerate
            num_list_match = re.match(r'^(\s*)\d+[.)]\s+(.+)', stripped)
            if num_list_match:
                depth = len(num_list_match.group(1)) // 2
                text = num_list_match.group(2).strip()
                runs = self._md_parse_inline(text)
                self.elements.append(DocumentElement(
                    element_type=ElementType.LIST_ITEM,
                    runs=runs,
                    list_type=ListType.NUMBERED,
                    list_depth=depth,
                ))
                in_list = True
                i += 1
                continue

            # Tabelle markdown
            if '|' in stripped and i + 1 < len(lines) and re.match(
                    r'^[\s|:-]+$', lines[i + 1].strip()):
                table_lines = []
                j = i
                while j < len(lines) and '|' in lines[j]:
                    table_lines.append(lines[j].strip())
                    j += 1

                rows = []
                for ti, tline in enumerate(table_lines):
                    if ti == 1 and re.match(r'^[\s|:-]+$', tline):
                        continue  # Riga separatore
                    cells_text = [c.strip() for c in tline.split('|')
                                  if c.strip() or ti == 0]
                    # Rimuovi celle vuote ai bordi
                    if cells_text and not cells_text[0]:
                        cells_text = cells_text[1:]
                    if cells_text and not cells_text[-1]:
                        cells_text = cells_text[:-1]
                    row_cells = [TableCell(runs=[TextRun(text=c)]) for c in cells_text]
                    if row_cells:
                        rows.append(row_cells)

                if rows:
                    self.elements.append(DocumentElement(
                        element_type=ElementType.TABLE,
                        table_rows=rows,
                        table_header_rows=1,
                    ))
                i = j
                continue

            # Paragrafo normale
            para_lines = [stripped]
            j = i + 1
            while j < len(lines):
                next_line = lines[j].rstrip('\n')
                if (not next_line.strip() or next_line.startswith('#')
                        or next_line.startswith('```')
                        or re.match(r'^[-*+]\s', next_line)
                        or re.match(r'^\d+[.)]\s', next_line)
                        or re.match(r'^[-*_]{3,}\s*$', next_line)):
                    break
                para_lines.append(next_line)
                j += 1

            text = ' '.join(l.strip() for l in para_lines)
            runs = self._md_parse_inline(text)
            self.elements.append(DocumentElement(
                element_type=ElementType.PARAGRAPH,
                runs=runs,
            ))
            i = j

    # =========================================================================
    # IPYNB PARSING
    # =========================================================================

    def _parse_ipynb(self):
        """Parsa un file Jupyter Notebook (.ipynb)."""
        with open(str(self.file_path), 'r', encoding='utf-8', errors='replace') as f:
            notebook = json.load(f)

        nb_meta = notebook.get('metadata', {})
        title = nb_meta.get('title') or self.file_path.stem
        self.metadata['title'] = title

        cells = notebook.get('cells', [])

        # Titolo notebook
        self.elements.append(DocumentElement(
            element_type=ElementType.HEADING,
            runs=[TextRun(text=title)],
            level=0,
        ))

        for cell in cells:
            cell_type = cell.get('cell_type', '')
            source = cell.get('source', '')
            if isinstance(source, list):
                source_text = ''.join(source)
            else:
                source_text = str(source or '')

            source_text = source_text.strip('\n')
            if not source_text.strip():
                continue

            if cell_type == 'markdown':
                self._parse_ipynb_markdown_cell(source_text)

            elif cell_type == 'code':
                self.elements.append(DocumentElement(
                    element_type=ElementType.CODE_BLOCK,
                    runs=[TextRun(text=source_text)],
                ))

                # Output testo del notebook (se presente)
                output_texts = []
                for output in cell.get('outputs', []):
                    text = ''
                    if 'text' in output:
                        data = output.get('text', '')
                        if isinstance(data, list):
                            text = ''.join(data)
                        else:
                            text = str(data)
                    elif 'data' in output:
                        plain = output.get('data', {}).get('text/plain', '')
                        if isinstance(plain, list):
                            text = ''.join(plain)
                        else:
                            text = str(plain)

                    text = text.strip('\n')
                    if text:
                        output_texts.append(text)

                if output_texts:
                    self.elements.append(DocumentElement(
                        element_type=ElementType.PARAGRAPH,
                        runs=[TextRun(text='Output:')],
                    ))
                    self.elements.append(DocumentElement(
                        element_type=ElementType.CODE_BLOCK,
                        runs=[TextRun(text='\n\n'.join(output_texts))],
                    ))

    def _parse_ipynb_markdown_cell(self, text: str):
        """Parsa il contenuto markdown di una cella notebook."""
        lines = text.splitlines()
        i = 0

        while i < len(lines):
            stripped = lines[i].rstrip('\n')

            if not stripped.strip():
                i += 1
                continue

            # Heading
            heading_match = re.match(r'^(#{1,6})\s+(.+)', stripped)
            if heading_match:
                level = len(heading_match.group(1))
                heading_text = heading_match.group(2).strip()
                self.elements.append(DocumentElement(
                    element_type=ElementType.HEADING,
                    runs=self._md_parse_inline(heading_text),
                    level=level,
                ))
                i += 1
                continue

            # Horizontal rule
            if re.match(r'^[-*_]{3,}\s*$', stripped):
                self.elements.append(DocumentElement(
                    element_type=ElementType.HORIZONTAL_RULE
                ))
                i += 1
                continue

            # Liste puntate
            list_match = re.match(r'^(\s*)[-*+]\s+(.+)', stripped)
            if list_match:
                depth = len(list_match.group(1)) // 2
                item_text = list_match.group(2).strip()
                self.elements.append(DocumentElement(
                    element_type=ElementType.LIST_ITEM,
                    runs=self._md_parse_inline(item_text),
                    list_type=ListType.BULLET,
                    list_depth=depth,
                ))
                i += 1
                continue

            # Liste numerate
            num_match = re.match(r'^(\s*)\d+[.)]\s+(.+)', stripped)
            if num_match:
                depth = len(num_match.group(1)) // 2
                item_text = num_match.group(2).strip()
                self.elements.append(DocumentElement(
                    element_type=ElementType.LIST_ITEM,
                    runs=self._md_parse_inline(item_text),
                    list_type=ListType.NUMBERED,
                    list_depth=depth,
                ))
                i += 1
                continue

            # Paragrafo
            para_lines = [stripped]
            j = i + 1
            while j < len(lines):
                next_line = lines[j].rstrip('\n')
                if (not next_line.strip() or next_line.startswith('#')
                        or re.match(r'^[-*+]\s', next_line)
                        or re.match(r'^\d+[.)]\s', next_line)
                        or re.match(r'^[-*_]{3,}\s*$', next_line)):
                    break
                para_lines.append(next_line)
                j += 1

            paragraph = ' '.join(line.strip() for line in para_lines)
            self.elements.append(DocumentElement(
                element_type=ElementType.PARAGRAPH,
                runs=self._md_parse_inline(paragraph),
            ))
            i = j

    def _md_parse_inline(self, text: str) -> list:
        """Parsa la formattazione inline del Markdown."""
        runs = []
        # Pattern: **bold**, *italic*, `code`, ~~strikethrough~~, [link](url)
        pattern = re.compile(
            r'(\*\*\*(.+?)\*\*\*'        # bold+italic
            r'|\*\*(.+?)\*\*'             # bold
            r'|\*(.+?)\*'                 # italic
            r'|`(.+?)`'                   # code
            r'|~~(.+?)~~'                 # strikethrough
            r'|\[([^\]]+)\]\(([^)]+)\)'   # link
            r')'
        )

        last_end = 0
        for match in pattern.finditer(text):
            # Testo prima del match
            if match.start() > last_end:
                runs.append(TextRun(text=text[last_end:match.start()]))

            if match.group(2):  # bold+italic
                runs.append(TextRun(
                    text=match.group(2),
                    styles={TextStyle.BOLD, TextStyle.ITALIC},
                ))
            elif match.group(3):  # bold
                runs.append(TextRun(
                    text=match.group(3),
                    styles={TextStyle.BOLD},
                ))
            elif match.group(4):  # italic
                runs.append(TextRun(
                    text=match.group(4),
                    styles={TextStyle.ITALIC},
                ))
            elif match.group(5):  # code
                runs.append(TextRun(
                    text=match.group(5),
                    styles={TextStyle.MONOSPACE},
                ))
            elif match.group(6):  # strikethrough
                runs.append(TextRun(
                    text=match.group(6),
                    styles={TextStyle.STRIKETHROUGH},
                ))
            elif match.group(7):  # link
                runs.append(TextRun(
                    text=match.group(7),
                    hyperlink=match.group(8),
                    styles={TextStyle.UNDERLINE},
                ))

            last_end = match.end()

        # Testo rimanente
        if last_end < len(text):
            runs.append(TextRun(text=text[last_end:]))

        if not runs:
            runs.append(TextRun(text=text))

        return runs
